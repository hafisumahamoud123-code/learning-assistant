import streamlit as st
import re
import json
import hashlib
import time
import random
from datetime import datetime, timedelta
from collections import Counter, defaultdict
import base64
from io import BytesIO
import os

# ========== FILE UPLOAD SUPPORT ==========
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Optional voice/OCR
try:
    from streamlit_webrtc import webrtc_streamer
    VOICE_AVAILABLE = True
except ImportError:
    VOICE_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ========== PAGE CONFIG ==========
st.set_page_config(
    page_title="Hafisu's AI Learning Companion",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ========== CUSTOM CSS ==========
st.markdown("""
<style>
    .main { background-color: #f0f2f6; }
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 20px;
        padding: 0.5rem 1rem;
        font-weight: bold;
    }
    .success-card {
        background: #d4edda;
        padding: 1rem;
        border-radius: 15px;
        margin: 0.5rem 0;
    }
    .warning-card {
        background: #fff3cd;
        padding: 1rem;
        border-radius: 15px;
        margin: 0.5rem 0;
    }
    .badge {
        display: inline-block;
        background: #667eea;
        color: white;
        border-radius: 30px;
        padding: 0.2rem 0.6rem;
        font-size: 0.7rem;
        margin-right: 0.5rem;
    }
    h1 { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
         -webkit-background-clip: text;
         -webkit-text-fill-color: transparent; }
</style>
""", unsafe_allow_html=True)

# ========== SESSION STATE INIT ==========
if 'user' not in st.session_state:
    st.session_state.user = {
        'name': 'Student',
        'xp': 0,
        'level': 1,
        'streak': 1,
        'last_active': str(datetime.now().date()),
        'subjects': defaultdict(int),
        'questions_answered': 0,
        'badges': ['🌟 First Login']
    }
if 'current_room' not in st.session_state:
    st.session_state.current_room = None
if 'room_messages' not in st.session_state:
    st.session_state.room_messages = []
if 'quiz_active' not in st.session_state:
    st.session_state.quiz_active = False
if 'quiz_questions' not in st.session_state:
    st.session_state.quiz_questions = []
if 'quiz_score' not in st.session_state:
    st.session_state.quiz_score = 0
if 'language' not in st.session_state:
    st.session_state.language = 'English'
if 'uploaded_text' not in st.session_state:
    st.session_state.uploaded_text = ""

# ========== HELPER FUNCTIONS ==========
def extract_text_from_file(uploaded_file):
    """Extract text from PDF, DOCX, PPTX, or TXT"""
    filename = uploaded_file.name.lower()
    content = uploaded_file.read()
    
    if filename.endswith('.txt'):
        return content.decode('utf-8', errors='ignore')
    
    elif filename.endswith('.pdf') and PDF_AVAILABLE:
        try:
            pdf = PyPDF2.PdfReader(BytesIO(content))
            text = ""
            for page in pdf.pages:
                text += page.extract_text()
            return text
        except:
            return "Error reading PDF. Make sure it's not password protected."
    
    elif filename.endswith('.docx') and DOCX_AVAILABLE:
        try:
            doc = Document(BytesIO(content))
            text = "\n".join([p.text for p in doc.paragraphs])
            return text
        except:
            return "Error reading Word document."
    
    elif filename.endswith('.pptx') and PPTX_AVAILABLE:
        try:
            prs = Presentation(BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except:
            return "Error reading PowerPoint file."
    
    else:
        return "Unsupported file type or missing library. Install: pip install PyPDF2 python-docx python-pptx"

def update_streak():
    today = str(datetime.now().date())
    last = st.session_state.user['last_active']
    if last != today:
        yesterday = (datetime.now() - timedelta(days=1)).date()
        if last == str(yesterday):
            st.session_state.user['streak'] += 1
        else:
            st.session_state.user['streak'] = 1
        st.session_state.user['last_active'] = today
        add_xp(10, reason="Daily login")

def add_xp(amount, reason="learning"):
    st.session_state.user['xp'] += amount
    new_level = st.session_state.user['xp'] // 100 + 1
    if new_level > st.session_state.user['level']:
        st.session_state.user['level'] = new_level
        st.session_state.user['badges'].append(f"🏆 Level {new_level}!")
    if st.session_state.user['streak'] >= 7 and "🔥 7 Day Streak" not in st.session_state.user['badges']:
        st.session_state.user['badges'].append("🔥 7 Day Streak")
    save_user_data()

def save_user_data():
    with open("user_data.json", "w") as f:
        data = dict(st.session_state.user)
        data['subjects'] = dict(data['subjects'])
        json.dump(data, f)

def load_user_data():
    try:
        with open("user_data.json", "r") as f:
            data = json.load(f)
            data['subjects'] = defaultdict(int, data.get('subjects', {}))
            st.session_state.user = data
    except:
        pass

def smart_summary(text, max_sentences=5):
    if len(text) < 50:
        return "Too short to summarize."
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if len(s.split()) > 5]
    if len(sentences) <= max_sentences:
        return text
    summary = ". ".join(sentences[:2] + sentences[-2:]) + "."
    return summary

def generate_flashcards(text, num=5):
    words = re.findall(r'\b[A-Za-z]{4,}\b', text.lower())
    stop = {'the','and','for','are','with','this','that','from','have','will','was','were','their','they'}
    words = [w for w in words if w not in stop]
    freq = Counter(words)
    top = [w for w,c in freq.most_common(num) if c>1]
    flashcards = []
    for w in top:
        flashcards.append({"term": w.capitalize(), "definition": f"Explain the concept of {w} based on the text."})
    return flashcards

def adaptive_quiz(text, difficulty):
    questions = []
    if difficulty == 'easy':
        terms = re.findall(r'\b[A-Z]{4,}\b', text) or re.findall(r'\b[a-z]{5,}\b', text)
        terms = list(set([t.capitalize() for t in terms[:5]]))
        for t in terms:
            questions.append({"question": f"What is {t}?", "options": ["A definition from the text", "A random fact", "None of the above"], "answer": 0})
    elif difficulty == 'medium':
        questions.append({"question": "What is the main argument?", "options": ["First sentence", "Last sentence", "Central idea"], "answer": 2})
    else:
        questions.append({"question": "How could you apply this knowledge in real life?", "options": ["Give an example", "Summarize again", "Ignore it"], "answer": 0})
    return questions

# ========== SIDEBAR PROFILE ==========
with st.sidebar:
    st.image("https://img.icons8.com/fluency/96/graduation-cap.png", width=80)
    st.markdown(f"## 👤 {st.session_state.user['name']}")
    col1, col2 = st.columns(2)
    col1.metric("⭐ XP", st.session_state.user['xp'])
    col2.metric("📊 Level", st.session_state.user['level'])
    st.progress(st.session_state.user['xp'] % 100 / 100, text="Next level progress")
    st.metric("🔥 Streak", f"{st.session_state.user['streak']} days")
    st.markdown("### 🏅 Badges")
    for badge in st.session_state.user['badges'][-5:]:
        st.markdown(f"- {badge}")
    st.markdown("---")
    st.session_state.language = st.selectbox("🌍 Language", ["English", "Twi", "Ga", "Hausa", "Fante", "French"])
    st.markdown("---")
    if st.button("🔄 Sync Data", use_container_width=True):
        load_user_data()
        st.success("Data reloaded!")

# ========== MAIN DASHBOARD ==========
st.title("🧠 Hafisu's AI Learning Companion")
st.markdown("#### *Created by Hafisu Mahamoud – BSc Agriculture, University of Ghana*")
st.markdown("**The world's first voice‑enabled, gamified, collaborative AI tutor for African students.**")

update_streak()

# Create tabs – added "About Creator" and "Upload File"
tabs = st.tabs([
    "📝 Summarize & Cards", 
    "🎯 Adaptive Quiz", 
    "💬 Collaborative Room", 
    "🎤 Voice Assistant", 
    "📷 Scan Textbook", 
    "📂 Upload File", 
    "👨‍🌾 About Creator"
])

# ---------- TAB 1: SUMMARIZE & FLASHCARDS ----------
with tabs[0]:
    colA, colB = st.columns([3,2])
    with colA:
        text_input = st.text_area("📖 Paste your lecture notes / textbook section:", height=250)
        if st.button("✨ Generate Summary & Flashcards", use_container_width=True):
            if text_input:
                summary = smart_summary(text_input)
                flashcards = generate_flashcards(text_input)
                st.session_state.last_summary = summary
                st.session_state.last_flashcards = flashcards
                add_xp(15, reason="Summarized material")
                st.success("✅ Processed! Check the right column.")
            else:
                st.warning("Please paste some text.")
    with colB:
        if 'last_summary' in st.session_state:
            st.markdown("### 📄 Summary")
            st.info(st.session_state.last_summary)
            st.markdown("### 🃏 Flashcards")
            for i, card in enumerate(st.session_state.last_flashcards):
                with st.expander(f"📌 {card['term']}"):
                    st.write(card['definition'])
        else:
            st.info("Your summary and flashcards will appear here.")

# ---------- TAB 2: ADAPTIVE QUIZ ----------
with tabs[1]:
    if not st.session_state.quiz_active:
        quiz_text = st.text_area("Paste text to quiz yourself on:", height=150)
        diff = st.select_slider("Difficulty", options=["Easy", "Medium", "Hard"])
        if st.button("🎮 Start Adaptive Quiz", use_container_width=True):
            if quiz_text:
                qs = adaptive_quiz(quiz_text, diff.lower())
                st.session_state.quiz_questions = qs
                st.session_state.quiz_active = True
                st.session_state.quiz_score = 0
                st.session_state.quiz_index = 0
                st.rerun()
            else:
                st.warning("Please paste text first.")
    else:
        idx = st.session_state.quiz_index
        if idx < len(st.session_state.quiz_questions):
            q = st.session_state.quiz_questions[idx]
            st.markdown(f"### Question {idx+1}/{len(st.session_state.quiz_questions)}")
            st.write(q['question'])
            answer = st.radio("Choose your answer:", q['options'], key=f"q{idx}")
            if st.button("Submit Answer"):
                if q['options'].index(answer) == q['answer']:
                    st.success("✅ Correct!")
                    st.session_state.quiz_score += 1
                    add_xp(10, reason="Quiz correct")
                else:
                    st.error(f"❌ Wrong! Correct: {q['options'][q['answer']]}")
                st.session_state.quiz_index += 1
                if st.session_state.quiz_index >= len(st.session_state.quiz_questions):
                    st.balloons()
                    st.success(f"🎉 Quiz completed! Score: {st.session_state.quiz_score}/{len(st.session_state.quiz_questions)}")
                    st.session_state.quiz_active = False
                st.rerun()
        else:
            st.session_state.quiz_active = False
            st.rerun()

# ---------- TAB 3: COLLABORATIVE ROOM ----------
with tabs[2]:
    room_code = st.text_input("🔑 Enter room code (or create new):")
    if st.button("🚪 Join / Create Room"):
        st.session_state.current_room = room_code if room_code else f"room_{random.randint(1000,9999)}"
        st.success(f"You are now in room: {st.session_state.current_room}")
    if st.session_state.current_room:
        st.markdown(f"### 💬 Room: {st.session_state.current_room}")
        chat_input = st.text_area("Share a question or note:")
        if st.button("📢 Post to Room"):
            if chat_input:
                st.session_state.room_messages.append({
                    "user": st.session_state.user['name'],
                    "msg": chat_input,
                    "time": datetime.now().strftime("%H:%M")
                })
                add_xp(5, reason="Collaboration")
        st.markdown("### Messages")
        for msg in st.session_state.room_messages[-20:]:
            st.markdown(f"**{msg['user']}** [{msg['time']}]: {msg['msg']}")

# ---------- TAB 4: VOICE ASSISTANT ----------
with tabs[3]:
    st.markdown("🎤 **Speak in your language** – the AI will understand and respond.")
    if VOICE_AVAILABLE:
        webrtc_streamer(key="speech", media_stream_constraints={"audio": True})
        st.caption("After speaking, your text will appear below.")
    else:
        st.warning("For full voice features, install: `pip install streamlit-webrtc SpeechRecognition`")
    st.markdown("### 📝 Or type your question:")
    voice_q = st.text_input("Your question:")
    if voice_q:
        st.info("🧠 The AI is listening. (Full voice integration with Google Speech API coming soon)")

# ---------- TAB 5: SCAN TEXTBOOK (OCR) ----------
with tabs[4]:
    st.markdown("📸 **Take a photo of a textbook page – we'll extract the text!**")
    uploaded_img = st.file_uploader("Upload an image (JPG/PNG)", type=["jpg","png","jpeg"])
    if uploaded_img and OCR_AVAILABLE:
        img = Image.open(uploaded_img)
        st.image(img, caption="Uploaded page", width=300)
        extracted = pytesseract.image_to_string(img)
        st.text_area("Extracted Text:", extracted, height=200)
        if st.button("Summarize extracted text"):
            sum_ocr = smart_summary(extracted)
            st.success(sum_ocr)
            add_xp(20, reason="Scanned textbook")
    elif uploaded_img and not OCR_AVAILABLE:
        st.error("OCR not installed. Run: `pip install pytesseract pillow` and install Tesseract from https://github.com/tesseract-ocr/tesseract")
    else:
        st.info("Upload an image of your textbook page to convert to text and summarize.")

# ---------- TAB 6: UPLOAD FILE (SLIDES, NOTES, DOCUMENTS) ----------
with tabs[5]:
    st.markdown("## 📂 Upload Full Course Materials")
    st.markdown("Upload **PDF, PowerPoint, Word documents, or Text files** – our AI will read and summarize them.")
    
    uploaded_file = st.file_uploader("Choose a file", type=['pdf', 'docx', 'pptx', 'txt'])
    if uploaded_file is not None:
        with st.spinner("Extracting text from file..."):
            file_text = extract_text_from_file(uploaded_file)
            if file_text and not file_text.startswith("Error") and not file_text.startswith("Unsupported"):
                st.session_state.uploaded_text = file_text
                st.success(f"✅ Successfully extracted {len(file_text.split())} words!")
                
                # Show preview
                with st.expander("Preview extracted text"):
                    st.text(file_text[:1000] + ("..." if len(file_text) > 1000 else ""))
                
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("📖 Summarize this document", use_container_width=True):
                        summary = smart_summary(file_text)
                        st.markdown("### Summary")
                        st.info(summary)
                        add_xp(15, reason="Uploaded and summarized file")
                with col2:
                    if st.button("🃏 Generate Flashcards", use_container_width=True):
                        flashcards = generate_flashcards(file_text)
                        st.markdown("### Flashcards")
                        for card in flashcards:
                            with st.expander(card['term']):
                                st.write(card['definition'])
                        add_xp(10, reason="Generated flashcards")
            else:
                st.error(file_text or "Could not process file. Check file format or install required libraries.")
    
    # Option to use previously uploaded text
    if st.session_state.uploaded_text:
        st.markdown("---")
        st.markdown("### Current document in memory")
        st.caption(f"{len(st.session_state.uploaded_text.split())} words loaded. Use the buttons above to summarize or create flashcards.")

# ---------- TAB 7: ABOUT THE CREATOR ----------
with tabs[6]:
    st.markdown("""
    ## 👨‍🌾 Meet the Creator: Hafisu Mahamoud
    
    | | |
    |---|---|
    | **Name** | Hafisu Mahamoud |
    | **Course** | BSc Agriculture |
    | **University** | University of Ghana, Legon |
    | **Year** | 2024/2025 |
    
    ### 📖 My Story
    
    > "I grew up speaking Hausa and English. In lecture halls, I saw friends from rural areas struggle to understand complex terms in English, even though they knew the concepts in Twi, Ga, or Fante. That broke my heart.
    >
    > I built **Hafisu's AI Learning Companion** so every student – whether you speak English, Twi, Ga, Hausa, Fante, or French – can learn in the language your mind understands best.
    >
    > No more language barriers. No more feeling left behind. Just you, your mother tongue, and the boundless joy of learning."
    
    ### 🎯 My Mission
    
    **Break language barriers in education for ALL students across ALL subjects.**
    
    ### 💬 My Quote
    
    > *"Knowledge grows best in your mother tongue. Learn in your language, achieve your dreams."*
    
    ### 🌟 Why This App is Different
    
    - **Voice & Local Languages** – Speak Twi, Ga, Hausa, and the AI listens
    - **Collaborative Study Rooms** – Learn with friends, earn XP together
    - **Gamification** – Level up, earn badges, keep streaks alive
    - **File Upload** – Upload slides, notes, PDFs – we summarize everything
    - **OCR Scanner** – Take a photo of any textbook page, AI reads it for you
    
    ### 🤝 Connect & Contribute
    
    This AI is **100% free** for every student in Ghana and Africa.
    
    Share it with your classmates. Use it daily. Let's break language barriers together.
    
    ---
    *Built with ❤️ by Hafisu Mahamoud – University of Ghana, BSc Agriculture*
    """)
    
    st.image("https://img.icons8.com/color/96/ghana.png", width=60)
    st.caption("Proudly Ghanaian 🇬🇭 | For every student, in every language.")

# ---------- FOOTER ----------
st.markdown("---")
st.markdown("💡 **Pro tip:** Use the **Upload File** tab to upload your lecture slides or notes. The AI will extract, summarize, and create flashcards for you!")
st.caption("© 2025 Hafisu Mahamoud | University of Ghana, BSc Agriculture | Built for every student in Africa 🇬🇭")

# Load existing user data on startup
load_user_data()